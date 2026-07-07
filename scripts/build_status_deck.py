#!/usr/bin/env python3
"""Build the NutriTrack Status Report deck (.pptx).

Run: python3 scripts/build_status_deck.py
Output: docs/status-report-2026-07-08.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy
import os

# ── Brand colors ─────────────────────────────────────────────────────────
ORANGE       = RGBColor(0xFF, 0x6B, 0x35)
ORANGE_DARK  = RGBColor(0xE8, 0x5A, 0x2C)
ORANGE_LIGHT = RGBColor(0xFF, 0x8A, 0x5C)
AMBER        = RGBColor(0xFF, 0xB6, 0x27)
CREAM        = RGBColor(0xFF, 0xF6, 0xEE)
SURFACE      = RGBColor(0xFA, 0xFA, 0xFA)
SURFACE_MUTE = RGBColor(0xF5, 0xF5, 0xF7)
TEXT_PRI     = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_SEC     = RGBColor(0x6B, 0x6B, 0x6B)
TEXT_TER     = RGBColor(0xA0, 0xA0, 0xA0)
GREEN        = RGBColor(0x00, 0xC8, 0x96)
RED          = RGBColor(0xE5, 0x39, 0x35)
RED_DEEP     = RGBColor(0xC6, 0x28, 0x28)
YELLOW       = RGBColor(0xFF, 0xB6, 0x27)
SKY          = RGBColor(0x4F, 0xC3, 0xF7)
DIVIDER      = RGBColor(0xEA, 0xEA, 0xEA)

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT  = os.path.join(ROOT, 'docs', 'status-report-2026-07-08.pptx')

# 16:9 widescreen
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK = prs.slide_layouts[6]

# ── Helpers ─────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s

def add_round(slide, x, y, w, h, fill, line=None, radius=0.06):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.adjustments[0] = radius
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s

def add_text(slide, x, y, w, h, text, *,
             font='Inter', size=14, bold=False, color=TEXT_PRI,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        lines = text.split('\n')
    else:
        lines = text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb

def add_rich(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """runs = [(text, dict_of_attrs)] per line. attrs: size, bold, color, font."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line_runs in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for text, attrs in line_runs:
            r = p.add_run()
            r.text = text
            r.font.name = attrs.get('font', 'Inter')
            r.font.size = Pt(attrs.get('size', 14))
            r.font.bold = attrs.get('bold', False)
            r.font.color.rgb = attrs.get('color', TEXT_PRI)
    return tb

def page_bg(slide, color=SURFACE):
    bg = add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, color)
    bg.shadow.inherit = False
    return bg

def page_header(slide, kicker, title, subtitle=None):
    # Top brand bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), ORANGE)
    # Kicker
    add_text(slide, Inches(0.5), Inches(0.25), Inches(6), Inches(0.3),
             kicker.upper(), size=10, bold=True, color=ORANGE,
             font='Inter')
    # Title
    add_text(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.55),
             title, size=28, bold=True, color=TEXT_PRI, font='Inter')
    # Subtitle (optional)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.05), Inches(12), Inches(0.4),
                 subtitle, size=14, color=TEXT_SEC, font='Inter')

def page_footer(slide, n, total):
    add_text(slide, Inches(0.5), Inches(7.05), Inches(6), Inches(0.3),
             "NutriTrack · Status Report · 2026-07-08", size=9, color=TEXT_TER, font='Inter')
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
             f"{n} / {total}", size=9, color=TEXT_TER, font='Inter', align=PP_ALIGN.RIGHT)

# ── Slide 1: Cover ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
page_bg(s, SURFACE)
# Big orange band on left
add_rect(s, 0, 0, Inches(4.5), SLIDE_H, ORANGE)
# Inner darker wedge (decorative)
add_rect(s, Inches(3.8), 0, Inches(0.7), SLIDE_H, ORANGE_DARK)
# Big "NT" mark on orange band — abstract letters
add_text(s, Inches(0.4), Inches(5.4), Inches(3.5), Inches(2),
         "NT", size=180, bold=True, color=CREAM, font='Inter',
         align=PP_ALIGN.LEFT)
add_text(s, Inches(0.5), Inches(0.5), Inches(3.5), Inches(0.4),
         "NUTRITRACK", size=11, bold=True, color=CREAM, font='Inter')
add_text(s, Inches(0.5), Inches(0.85), Inches(3.5), Inches(0.3),
         "v0.1.0  ·  build 2026.07.08", size=10, color=CREAM, font='Inter')

# Title block (right side)
add_text(s, Inches(5.0), Inches(1.2), Inches(8), Inches(0.4),
         "STATUS REPORT", size=14, bold=True, color=ORANGE, font='Inter')
add_text(s, Inches(5.0), Inches(1.65), Inches(8), Inches(1.6),
         "Platform, infrastructure,\narchitecture & UX", size=32, bold=True,
         color=TEXT_PRI, font='Inter', line_spacing=1.15)
add_text(s, Inches(5.0), Inches(3.6), Inches(8), Inches(0.5),
         "Top-priority problem list with risk & urgency", size=18, color=TEXT_SEC, font='Inter')

# Date / author block
add_text(s, Inches(5.0), Inches(5.5), Inches(8), Inches(0.3),
         "REPORT DATE", size=9, bold=True, color=TEXT_TER, font='Inter')
add_text(s, Inches(5.0), Inches(5.8), Inches(8), Inches(0.4),
         "2026-07-08  ·  00:12 Asia/Shanghai", size=14, color=TEXT_PRI, font='Inter')
add_text(s, Inches(5.0), Inches(6.3), Inches(8), Inches(0.3),
         "AUTHOR", size=9, bold=True, color=TEXT_TER, font='Inter')
add_text(s, Inches(5.0), Inches(6.55), Inches(8), Inches(0.4),
         "Mavis (audit + analysis)", size=14, color=TEXT_PRI, font='Inter')

# ── Slide 2: Executive Summary ──────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
page_bg(s)
page_header(s, "Section 01", "Executive Summary",
           "Where NutriTrack stands today, in 4 numbers")

# 4 KPI tiles
tiles = [
    ("30%",  "Production-ready", RED, "Code is structured\nbut doesn't build"),
    ("9,373","Dart LOC",        TEXT_PRI, "across 6 features,\nclean architecture"),
    ("31",   "Compile errors",   RED, "block next build,\nall in 2 files"),
    ("0",    "Tests",            RED, "zero coverage,\ntest/ doesn't exist"),
]
x = Inches(0.5); y = Inches(1.6); w = Inches(3.0); h = Inches(2.6); gap = Inches(0.15)
for i, (big, label, color, sub) in enumerate(tiles):
    add_round(s, x + i*(w+gap), y, w, h, SURFACE_MUTE, line=DIVIDER)
    add_text(s, x + i*(w+gap), y + Inches(0.4), w, Inches(1.0),
             big, size=48, bold=True, color=color, font='Inter',
             align=PP_ALIGN.CENTER)
    add_text(s, x + i*(w+gap), y + Inches(1.5), w, Inches(0.4),
             label, size=14, bold=True, color=TEXT_PRI, font='Inter',
             align=PP_ALIGN.CENTER)
    add_text(s, x + i*(w+gap) + Inches(0.3), y + Inches(1.95), w - Inches(0.6), Inches(0.7),
             sub, size=10, color=TEXT_SEC, font='Inter',
             align=PP_ALIGN.CENTER, line_spacing=1.3)

# 3 takeaway boxes
y2 = Inches(4.55); h2 = Inches(2.3)
add_round(s, Inches(0.5), y2, Inches(4.05), h2, SURFACE, line=DIVIDER)
add_rect(s, Inches(0.5), y2, Inches(0.12), h2, GREEN)
add_text(s, Inches(0.85), y2 + Inches(0.2), Inches(3.5), Inches(0.4),
         "WHAT'S WORKING", size=10, bold=True, color=GREEN, font='Inter')
add_rich(s, Inches(0.85), y2 + Inches(0.6), Inches(3.5), Inches(1.8), [
    [("✓ ", {'size':14,'bold':True,'color':GREEN}),
     ("AI gateway architecture\n", {'size':12,'color':TEXT_PRI})],
    [("✓ ", {'size':14,'bold':True,'color':GREEN}),
     ("Drift DB schema (7 tables)\n", {'size':12,'color':TEXT_PRI})],
    [("✓ ", {'size':14,'bold':True,'color':GREEN}),
     ("PB live (8 nt_* collections)\n", {'size':12,'color':TEXT_PRI})],
    [("✓ ", {'size':14,'bold':True,'color':GREEN}),
     ("Brand identity shipped\n", {'size':12,'color':TEXT_PRI})],
    [("✓ ", {'size':14,'bold':True,'color':GREEN}),
     ("AI vision + offline cache", {'size':12,'color':TEXT_PRI})],
], line_spacing=1.4)

add_round(s, Inches(4.7), y2, Inches(4.05), h2, SURFACE, line=DIVIDER)
add_rect(s, Inches(4.7), y2, Inches(0.12), h2, YELLOW)
add_text(s, Inches(5.05), y2 + Inches(0.2), Inches(3.5), Inches(0.4),
         "HALF-BUILT", size=10, bold=True, color=YELLOW, font='Inter')
add_rich(s, Inches(5.05), y2 + Inches(0.6), Inches(3.5), Inches(1.8), [
    [("⚠ ", {'size':14,'bold':True,'color':YELLOW}),
     ("5/6 screens render fake data\n", {'size':12,'color':TEXT_PRI})],
    [("⚠ ", {'size':14,'bold':True,'color':YELLOW}),
     ("No auth / onboarding route\n", {'size':12,'color':TEXT_PRI})],
    [("⚠ ", {'size':14,'bold':True,'color':YELLOW}),
     ("Workout session UI missing\n", {'size':12,'color':TEXT_PRI})],
    [("⚠ ", {'size':14,'bold':True,'color':YELLOW}),
     ("Exercise DB: 46/488 (9.4%)\n", {'size':12,'color':TEXT_PRI})],
    [("⚠ ", {'size':14,'bold':True,'color':YELLOW}),
     ("No Sentry / crash reporting", {'size':12,'color':TEXT_PRI})],
], line_spacing=1.4)

add_round(s, Inches(8.9), y2, Inches(3.9), h2, SURFACE, line=DIVIDER)
add_rect(s, Inches(8.9), y2, Inches(0.12), h2, RED)
add_text(s, Inches(9.25), y2 + Inches(0.2), Inches(3.5), Inches(0.4),
         "BLOCKERS", size=10, bold=True, color=RED, font='Inter')
add_rich(s, Inches(9.25), y2 + Inches(0.6), Inches(3.5), Inches(1.8), [
    [("✗ ", {'size':14,'bold':True,'color':RED}),
     ("App won't compile\n", {'size':12,'color':TEXT_PRI})],
    [("✗ ", {'size':14,'bold':True,'color':RED}),
     ("Never run on a device\n", {'size':12,'color':TEXT_PRI})],
    [("✗ ", {'size':14,'bold':True,'color':RED}),
     ("9 packages on major N-1\n", {'size':12,'color':TEXT_PRI})],
    [("✗ ", {'size':14,'bold':True,'color':RED}),
     ("No crash reporting\n", {'size':12,'color':TEXT_PRI})],
    [("✗ ", {'size':14,'bold':True,'color':RED}),
     ("7 dead dependencies", {'size':12,'color':TEXT_PRI})],
], line_spacing=1.4)

page_footer(s, 2, 14)


# ── Slide 3: Platform Status ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
page_bg(s)
page_header(s, "Section 02", "Platform Status",
           "Flutter SDK, dependencies, platform readiness")

# 2 columns: stack + versions
add_text(s, Inches(0.5), Inches(1.6), Inches(6), Inches(0.4),
         "STACK SNAPSHOT", size=10, bold=True, color=ORANGE, font='Inter')

stack = [
    ("Flutter",        "3.27+ (declared) · 3.29.3 in your lockfile", "✓"),
    ("Dart",           "3.6+ · supports records, patterns, sealed",  "✓"),
    ("State",          "Riverpod 2.6.1 with @riverpod codegen",       "✓"),
    ("Routing",        "go_router 14.8.1 (StatefulShellRoute)",       "✓"),
    ("Local DB",       "Drift 2.28.2 (7 tables, all defined)",        "✓"),
    ("AI Gateway",     "OpenRouter → MiniMax M3 + Gemini + GPT-4o", "✓"),
    ("Voice",          "record 5.1.2 → Whisper transcription",        "✓"),
    ("Barcode",        "mobile_scanner 5.2.3 + OFF client",          "⚠"),
    ("Camera",         "camera 0.11.4 → 0.12.0+1 available",        "⚠"),
    ("Cloud sync",     "Supabase declared but unused",               "⚠"),
    ("Health",         "REMOVED in this sprint (incompat)",          "✗"),
]
y = Inches(2.0)
for i, (k, v, status) in enumerate(stack):
    color = GREEN if status == '✓' else (YELLOW if status == '⚠' else RED)
    add_round(s, Inches(0.5), y + i*Inches(0.36), Inches(6.3), Inches(0.32),
              SURFACE if i % 2 == 0 else SURFACE_MUTE, line=DIVIDER, radius=0.15)
    add_text(s, Inches(0.7), y + i*Inches(0.36) + Inches(0.04), Inches(1.4), Inches(0.28),
             k, size=11, bold=True, color=TEXT_PRI, font='Inter')
    add_text(s, Inches(2.1), y + i*Inches(0.36) + Inches(0.04), Inches(4.2), Inches(0.28),
             v, size=10, color=TEXT_SEC, font='Inter')
    add_text(s, Inches(6.4), y + i*Inches(0.36) + Inches(0.04), Inches(0.3), Inches(0.28),
             status, size=14, bold=True, color=color, font='Inter')

# Right side: build readiness
add_text(s, Inches(7.0), Inches(1.6), Inches(6), Inches(0.4),
         "BUILD READINESS", size=10, bold=True, color=ORANGE, font='Inter')

readiness = [
    ("Android compile",    "✗ FAILS",  RED,    "31 errors in barcode/camera"),
    ("Android (post-fix)", "STAGING",  YELLOW, "~45 min of mechanical fixes"),
    ("iOS compile",        "UNKNOWN",  YELLOW, "macOS runner; not validated"),
    ("Device run",         "NEVER",    RED,    "No `flutter run` ever executed"),
    ("flutter analyze",    "31 errors + 17 warnings + 62 info", RED, "in 6 files"),
    ("dart fix --apply",   "WOULD KILL 12 of 17 warnings", GREEN, "5 sec cleanup"),
    ("pubspec.lock",       "GITIGNORED", YELLOW, "fresh clone = `pub get`"),
    ("Generated files",    "GITIGNORED", YELLOW, "fresh clone = `build_runner`"),
    ("Secrets config",     "TEMPLATE ONLY", YELLOW, "user must cp template"),
]
y = Inches(2.0)
for i, (k, v, c, sub) in enumerate(readiness):
    add_round(s, Inches(7.0), y + i*Inches(0.36), Inches(5.8), Inches(0.32),
              SURFACE if i % 2 == 0 else SURFACE_MUTE, line=DIVIDER, radius=0.15)
    add_text(s, Inches(7.2), y + i*Inches(0.36) + Inches(0.04), Inches(2.0), Inches(0.28),
             k, size=11, bold=True, color=TEXT_PRI, font='Inter')
    add_text(s, Inches(9.2), y + i*Inches(0.36) + Inches(0.04), Inches(1.6), Inches(0.28),
             v, size=10, bold=True, color=c, font='Inter')
    add_text(s, Inches(7.2), y + i*Inches(0.36) + Inches(0.22), Inches(5.4), Inches(0.18),
             sub, size=8, color=TEXT_TER, font='Inter')

page_footer(s, 3, 14)


# ── Slide 4: Infrastructure Status ──────────────────────────────────────
s = prs.slides.add_slide(BLANK)
page_bg(s)
page_header(s, "Section 03", "Infrastructure Status",
           "PocketBase, GitHub Actions, dependencies")

# 4 panels: PB schema, PB data, CI workflows, secrets
panel_y = Inches(1.7); panel_h = Inches(5.0)
panel_w = Inches(6.05); gap = Inches(0.2)

# Panel A: PB Schema
add_round(s, Inches(0.5), panel_y, panel_w, panel_h, SURFACE, line=DIVIDER)
add_rect(s, Inches(0.5), panel_y, Inches(0.12), panel_h, GREEN)
add_text(s, Inches(0.85), panel_y + Inches(0.2), Inches(5), Inches(0.4),
         "POCKETBASE · LIVE SCHEMA", size=10, bold=True, color=GREEN, font='Inter')
add_text(s, Inches(0.85), panel_y + Inches(0.55), Inches(5), Inches(0.4),
         "9 collections on pocketbase.scaleupcrm.com", size=12, bold=True, color=TEXT_PRI, font='Inter')

collections = [
    ("nt_users",           "0 records",  "auth-anchored profiles"),
    ("nt_food_logs",       "0 records",  "synced meal entries"),
    ("nt_exercises",       "46 records", "chest-only seeded"),
    ("nt_workout_sessions","0 records",  "session history"),
    ("nt_weight_entries",  "0 records",  "weight log"),
    ("nt_favorites",       "0 records",  "starred foods"),
    ("nt_meal_templates",  "0 records",  "recurring meals"),
    ("nt_sync_queue",      "0 records",  "offline→cloud queue"),
    ("nt_barcode_cache",   "0 records",  "OFF cache (anon read)"),
]
y = panel_y + Inches(1.05)
for i, (n, cnt, sub) in enumerate(collections):
    add_text(s, Inches(0.85), y + i*Inches(0.36), Inches(2.4), Inches(0.28),
             n, size=11, bold=True, color=TEXT_PRI, font='Inter')
    cnt_color = ORANGE if cnt != "0 records" else TEXT_TER
    add_text(s, Inches(3.25), y + i*Inches(0.36), Inches(1.2), Inches(0.28),
             cnt, size=10, color=cnt_color, font='Inter')
    add_text(s, Inches(4.45), y + i*Inches(0.36), Inches(2.0), Inches(0.28),
             sub, size=10, color=TEXT_SEC, font='Inter')

# Panel B: CI Workflows
add_round(s, Inches(6.85), panel_y, panel_w, panel_h, SURFACE, line=DIVIDER)
add_rect(s, Inches(6.85), panel_y, Inches(0.12), panel_h, ORANGE)
add_text(s, Inches(7.2), panel_y + Inches(0.2), Inches(5), Inches(0.4),
         "GITHUB ACTIONS · CI/CD", size=10, bold=True, color=ORANGE, font='Inter')
add_text(s, Inches(7.2), panel_y + Inches(0.55), Inches(5), Inches(0.4),
         "3 workflows + 1 README", size=12, bold=True, color=TEXT_PRI, font='Inter')

workflows = [
    ("ci.yml",            "ubuntu/macos runners · analyze + test + build APK/iOS", GREEN),
    ("off-to-pb-sync.yml","cron 0 */6h · fetches OFF products → PB cache",         GREEN),
    ("dependency-sync.yml","weekly · opens PR for outdated packages",             GREEN),
]
y = panel_y + Inches(1.05)
for i, (n, sub, c) in enumerate(workflows):
    add_text(s, Inches(7.2), y + i*Inches(0.7), Inches(5.5), Inches(0.3),
             n, size=12, bold=True, color=TEXT_PRI, font='Inter')
    add_text(s, Inches(7.2), y + i*Inches(0.7) + Inches(0.28), Inches(5.5), Inches(0.4),
             sub, size=10, color=TEXT_SEC, font='Inter')

add_text(s, Inches(7.2), panel_y + Inches(3.5), Inches(5), Inches(0.3),
         "SECRETS", size=10, bold=True, color=ORANGE, font='Inter')
secrets = [
    ("PB_URL",              "in workflow env"),
    ("PB_ADMIN_EMAIL",      "in workflow secrets"),
    ("PB_ADMIN_PASSWORD",   "in workflow secrets"),
    ("OPENROUTER_API_KEY",  "in secrets.dart (template only)"),
]
y = panel_y + Inches(3.85)
for i, (k, v) in enumerate(secrets):
    add_text(s, Inches(7.2), y + i*0.32*914400, Inches(2.0), Inches(0.28),
             k, size=11, bold=True, color=TEXT_PRI, font='Inter')
    add_text(s, Inches(9.2), y + i*0.32*914400, Inches(3.4), Inches(0.28),
             v, size=10, color=TEXT_SEC, font='Inter')

page_footer(s, 4, 14)


# ── Slide 5: Architecture ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
page_bg(s)
page_header(s, "Section 04", "Architecture",
           "Domain-Driven Design · 6 features · 7 Drift tables")

# Left: feature LOC table
add_text(s, Inches(0.5), Inches(1.6), Inches(6), Inches(0.4),
         "FEATURES BY SIZE", size=10, bold=True, color=ORANGE, font='Inter')

features = [
    ("camera",      1962, "snap → AI → review"),
    ("dashboard",   1969, "today's meals + macros"),
    ("barcode",     1654, "scan → OFF lookup"),
    ("workout",     1046, "search + log"),
    ("settings",     736, "TDEE + profile"),
    ("insights",     481, "weight + trends"),
]
y = Inches(2.0)
for i, (name, loc, sub) in enumerate(features):
    bar_w = loc / 25  # scale to ~80px max
    add_round(s, Inches(0.5), y + i*Inches(0.5), Inches(6.3), Inches(0.45),
              SURFACE if i % 2 == 0 else SURFACE_MUTE, line=DIVIDER, radius=0.2)
    add_text(s, Inches(0.7), y + i*Inches(0.5) + Inches(0.08), Inches(1.2), Inches(0.3),
             name, size=12, bold=True, color=TEXT_PRI, font='Inter')
    add_text(s, Inches(1.9), y + i*Inches(0.5) + Inches(0.08), Inches(4.0), Inches(0.3),
             sub, size=10, color=TEXT_SEC, font='Inter')
    # bar
    add_round(s, Inches(5.5), y + i*Inches(0.5) + Inches(0.13), Inches(int(bar_w*14400)), Inches(0.18),
              ORANGE if loc > 1000 else ORANGE_LIGHT)
    add_text(s, Inches(6.7), y + i*Inches(0.5) + Inches(0.08), Inches(0.4), Inches(0.3),
             f"{loc}", size=11, bold=True, color=TEXT_PRI, font='Inter')

# Right: Drift tables
add_text(s, Inches(7.0), Inches(1.6), Inches(6), Inches(0.4),
         "DRIFT TABLES (7)", size=10, bold=True, color=ORANGE, font='Inter')

tables = [
    ("FoodLogEntries",    "meal entries + macros + source"),
    ("ExerciseEntries",   "logged exercises per session"),
    ("WorkoutSessions",   "active + history sessions"),
    ("WeightEntries",     "weight log + date"),
    ("UserProfiles",      "TDEE + goal + biometrics"),
    ("ImageHashCache",    "AI dedup (avoid re-recognize)"),
    ("PendingSyncEntries","offline → cloud queue"),
]
y = Inches(2.0)
for i, (n, sub) in enumerate(tables):
    add_round(s, Inches(7.0), y + i*Inches(0.55), Inches(5.8), Inches(0.5),
              SURFACE if i % 2 == 0 else SURFACE_MUTE, line=DIVIDER, radius=0.15)
    add_text(s, Inches(7.2), y + i*Inches(0.55) + Inches(0.08), Inches(2.6), Inches(0.3),
             n, size=12, bold=True, color=TEXT_PRI, font='Inter')
    add_text(s, Inches(7.2), y + i*Inches(0.55) + Inches(0.28), Inches(5.0), Inches(0.25),
             sub, size=9, color=TEXT_SEC, font='Inter')

# Bottom strip: data flow
y3 = Inches(5.95)
add_round(s, Inches(0.5), y3, Inches(12.3), Inches(0.95), SURFACE_MUTE, line=DIVIDER)
add_text(s, Inches(0.7), y3 + Inches(0.1), Inches(11), Inches(0.3),
         "DATA FLOW", size=10, bold=True, color=ORANGE, font='Inter')
add_rich(s, Inches(0.7), y3 + Inches(0.4), Inches(12), Inches(0.5), [
    [("Capture", {'size':12,'bold':True,'color':TEXT_PRI}),
     (" → AI Gateway → ", {'size':12,'color':TEXT_SEC}),
     ("Drift (local)", {'size':12,'bold':True,'color':ORANGE}),
     (" → ", {'size':12,'color':TEXT_SEC}),
     ("PocketBase (cloud)", {'size':12,'bold':True,'color':GREEN}),
     ("  ·  3-tier barcode cache: in-memory → PB → OFF (anonymous read)", {'size':11,'color':TEXT_SEC})]
])

page_footer(s, 5, 14)


# ── Slide 6: Task List — Current Sprint ─────────────────────────────────
s = prs.slides.add_slide(BLANK)
page_bg(s)
page_header(s, "Section 05", "Current Task Status",
           "What's done, what's in progress, what's blocked")

# 3 columns
col_w = Inches(4.05); col_h = Inches(5.0); gap = Inches(0.15)
y = Inches(1.6)

# DONE
add_round(s, Inches(0.5), y, col_w, col_h, SURFACE, line=DIVIDER)
add_rect(s, Inches(0.5), y, Inches(0.12), Inches(0.5), GREEN)
add_text(s, Inches(0.85), y + Inches(0.1), Inches(2.5), Inches(0.4),
         "DONE", size=11, bold=True, color=GREEN, font='Inter')
add_text(s, Inches(3.3), y + Inches(0.1), Inches(1.0), Inches(0.4),
         "12", size=14, bold=True, color=GREEN, font='Inter', align=PP_ALIGN.RIGHT)
done = [
    "Drift migration (Isar removed)",
    "PB schema live (9 collections)",
    "46 chest exercises seeded",
    "AI gateway architecture",
    "Drift codegen working",
    "Brand identity (icon, palette)",
    "Brand exports (iOS/Android/Web)",
    "KNOWN_WARNINGS.md written",
    "Gradle 8.14.0 + AGP 8.11.1",
    "record 5.1.2 pin (Linux iface)",
    "health plugin removed",
    "media/ golden-source structure",
]
for i, item in enumerate(done):
    add_text(s, Inches(0.85), y + Inches(0.6) + i*Inches(0.34), Inches(0.2), Inches(0.3),
             "✓", size=12, bold=True, color=GREEN, font='Inter')
    add_text(s, Inches(1.05), y + Inches(0.6) + i*Inches(0.34), Inches(3.4), Inches(0.3),
             item, size=10, color=TEXT_PRI, font='Inter')

# IN PROGRESS
add_round(s, Inches(4.7), y, col_w, col_h, SURFACE, line=DIVIDER)
add_rect(s, Inches(4.7), y, Inches(0.12), Inches(0.5), ORANGE)
add_text(s, Inches(5.05), y + Inches(0.1), Inches(2.5), Inches(0.4),
         "IN PROGRESS", size=11, bold=True, color=ORANGE, font='Inter')
add_text(s, Inches(7.5), y + Inches(0.1), Inches(1.0), Inches(0.4),
         "3", size=14, bold=True, color=ORANGE, font='Inter', align=PP_ALIGN.RIGHT)
prog = [
    "Audit doc v2 (PRODUCTION_READINESS)",
    "Strategy doc (WHAT_MISSING)",
    "Major-version debt tracking (§6)",
]
for i, item in enumerate(prog):
    add_text(s, Inches(5.05), y + Inches(0.6) + i*Inches(0.5), Inches(0.2), Inches(0.3),
             "⟳", size=14, bold=True, color=ORANGE, font='Inter')
    add_rich(s, Inches(5.3), y + Inches(0.6) + i*Inches(0.5), Inches(3.4), Inches(0.45), [
        [(item, {'size':10,'color':TEXT_PRI})]
    ])

# BLOCKED
add_round(s, Inches(8.9), y, col_w, col_h, SURFACE, line=DIVIDER)
add_rect(s, Inches(8.9), y, Inches(0.12), Inches(0.5), RED)
add_text(s, Inches(9.25), y + Inches(0.1), Inches(2.5), Inches(0.4),
         "BLOCKED", size=11, bold=True, color=RED, font='Inter')
add_text(s, Inches(11.7), y + Inches(0.1), Inches(1.0), Inches(0.4),
         "5", size=14, bold=True, color=RED, font='Inter', align=PP_ALIGN.RIGHT)
blocked = [
    ("Android compile",          "31 errors in 2 files"),
    ("First device run",         "no working APK/IPA"),
    ("Insights screen",          "renders fake sample data"),
    ("Workout session UI",       "no active-session flow"),
    ("Onboarding flow",          "no /onboarding route"),
]
for i, (k, sub) in enumerate(blocked):
    add_text(s, Inches(9.25), y + Inches(0.6) + i*Inches(0.85), Inches(0.2), Inches(0.3),
             "✗", size=12, bold=True, color=RED, font='Inter')
    add_text(s, Inches(9.5), y + Inches(0.6) + i*Inches(0.85), Inches(3.0), Inches(0.3),
             k, size=11, bold=True, color=TEXT_PRI, font='Inter')
    add_text(s, Inches(9.5), y + Inches(0.6) + i*Inches(0.85) + Inches(0.28), Inches(3.0), Inches(0.5),
             sub, size=9, color=TEXT_SEC, font='Inter')

# Velocity strip at bottom
add_round(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5), SURFACE_MUTE, line=DIVIDER)
add_text(s, Inches(0.7), Inches(6.78), Inches(2), Inches(0.35),
         "VELOCITY", size=10, bold=True, color=ORANGE, font='Inter')
add_text(s, Inches(1.7), Inches(6.78), Inches(11), Inches(0.35),
         "9 commits in last 7 days · ~3 commits/day · clear blockers to unblock more",
         size=11, color=TEXT_PRI, font='Inter')

page_footer(s, 6, 14)


# ── Slide 7: UX Status ──────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
page_bg(s)
page_header(s, "Section 06", "UX Status",
           "What users would see today if the app compiled")

# 4 status categories
cats = [
    ("SCREENS", "6", "dashboard, workout,\ninsights, settings,\ncamera, barcode", GREEN),
    ("WIRED TO DATA", "1 of 6", "5 of 6 screens render\nplaceholder / sample data", RED),
    ("EMPTY STATES", "0", "no `shimmer`, no\nempty-state illustrations", RED),
    ("MOTION + HAPTIC", "minimal", "app_motion exists\nbut underused", YELLOW),
]
y = Inches(1.6); w = Inches(3.0); h = Inches(1.7); gap = Inches(0.15)
for i, (label, big, sub, c) in enumerate(cats):
    add_round(s, Inches(0.5) + i*(w+gap), y, w, h, SURFACE, line=DIVIDER)
    add_text(s, Inches(0.5) + i*(w+gap), y + Inches(0.15), w, Inches(0.3),
             label, size=10, bold=True, color=TEXT_TER, font='Inter', align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5) + i*(w+gap), y + Inches(0.45), w, Inches(0.6),
             big, size=36, bold=True, color=c, font='Inter', align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5) + i*(w+gap) + Inches(0.2), y + Inches(1.1), w - Inches(0.4), Inches(0.6),
             sub, size=10, color=TEXT_SEC, font='Inter', align=PP_ALIGN.CENTER, line_spacing=1.3)

# Per-screen status table
add_text(s, Inches(0.5), Inches(3.55), Inches(12), Inches(0.4),
         "SCREEN-BY-SCREEN UX READINESS", size=10, bold=True, color=ORANGE, font='Inter')
screens = [
    ("Dashboard",   "real data wiring", "fake data, no onboarding",      RED),
    ("Workout",     "search working",   "no session UI, no rest timer",  RED),
    ("Insights",    "fl_chart present", "renders `_generateSampleWeights()`", RED),
    ("Settings",    "TDEE wizard OK",   "no first-launch gate",          YELLOW),
    ("Camera",      "build broken",     "31 compile errors",             RED),
    ("Barcode",     "build broken",     "RRect.topLeft undefined",       RED),
]
y = Inches(4.0)
for i, (n, good, bad, c) in enumerate(screens):
    add_round(s, Inches(0.5), y + i*Inches(0.45), Inches(12.3), Inches(0.4),
              SURFACE if i % 2 == 0 else SURFACE_MUTE, line=DIVIDER, radius=0.2)
    add_text(s, Inches(0.7), y + i*Inches(0.45) + Inches(0.07), Inches(2.0), Inches(0.3),
             n, size=12, bold=True, color=TEXT_PRI, font='Inter')
    add_text(s, Inches(2.8), y + i*Inches(0.45) + Inches(0.07), Inches(4.5), Inches(0.3),
             "✓ " + good, size=10, color=GREEN, font='Inter')
    add_text(s, Inches(7.5), y + i*Inches(0.45) + Inches(0.07), Inches(4.5), Inches(0.3),
             "✗ " + bad, size=10, color=c, font='Inter')

page_footer(s, 7, 14)


# ── Slide 8: Performance Suggestions ────────────────────────────────────
s = prs.slides.add_slide(BLANK)
page_bg(s)
page_header(s, "Section 07", "Performance Suggestions",
           "What to fix before users notice the lag")

perfs = [
    ("1", "Skeleton loaders on every list",
     "Dashboard tiles, exercise search, history, insights charts all flash empty→full. Add shimmer widgets.",
     "~4 hours", "BIG", GREEN),
    ("2", "RepaintBoundary around MacroDonut",
     "Already added in earlier commit. Verify on real device — fl_chart redraws on every frame otherwise.",
     "~1 hour", "MED", YELLOW),
    ("3", "Cache Network images (cached_network_image)",
     "Barcode images already wrapped. Apply same pattern to food images from OFF / PB.",
     "~2 hours", "MED", YELLOW),
    ("4", "Lazy-load PB collections",
     "Currently FoodLogRepository fetches full day on every screen mount. Use Drift streams (watchByDate) instead of getByDate.",
     "~3 hours", "BIG", ORANGE),
    ("5", "Profile-mode build for production",
     "flutter build apk --release --obfuscate --split-debug-info=debug/. Cuts ~30% APK size.",
     "~30 min", "QUICK", GREEN),
    ("6", "Image compression before AI",
     "Camera capture = 2-4 MB JPEG. Resize to 1024px, JPEG quality 85 = ~200 KB. Faster AI upload, lower bandwidth.",
     "~2 hours", "MED", YELLOW),
    ("7", "Skia caching for static content",
     "macro_donut, wordmark, empty states — wrap in RepaintBoundary or const where possible.",
     "~2 hours", "MED", YELLOW),
    ("8", "Sync TextEditingController (already done)",
     "Workout search field — verify no sync creation in build(). You fixed this already. ✓",
     "✓", "DONE", GREEN),
]
y = Inches(1.55)
for i, (n, title, sub, effort, size, c) in enumerate(perfs):
    add_round(s, Inches(0.5), y + i*Inches(0.65), Inches(12.3), Inches(0.6),
              SURFACE if i % 2 == 0 else SURFACE_MUTE, line=DIVIDER, radius=0.2)
    # Number badge
    add_round(s, Inches(0.65), y + i*Inches(0.65) + Inches(0.1), Inches(0.4), Inches(0.4),
              c, radius=0.5)
    add_text(s, Inches(0.65), y + i*Inches(0.65) + Inches(0.13), Inches(0.4), Inches(0.35),
             n, size=14, bold=True, color=CREAM, font='Inter', align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.2), y + i*Inches(0.65) + Inches(0.05), Inches(7.5), Inches(0.3),
             title, size=12, bold=True, color=TEXT_PRI, font='Inter')
    add_text(s, Inches(1.2), y + i*Inches(0.65) + Inches(0.28), Inches(7.5), Inches(0.35),
             sub, size=9, color=TEXT_SEC, font='Inter', line_spacing=1.2)
    add_text(s, Inches(9.0), y + i*Inches(0.65) + Inches(0.1), Inches(2.5), Inches(0.3),
             effort, size=11, bold=True, color=TEXT_PRI, font='Inter', align=PP_ALIGN.CENTER)
    add_text(s, Inches(11.5), y + i*Inches(0.65) + Inches(0.1), Inches(1.0), Inches(0.3),
             size, size=10, bold=True, color=c, font='Inter', align=PP_ALIGN.CENTER)

page_footer(s, 8, 14)


# ── Slide 9: Security Suggestions ───────────────────────────────────────
s = prs.slides.add_slide(BLANK)
page_bg(s)
page_header(s, "Section 08", "Security Suggestions",
           "What to lock down before the first user logs in")

sec = [
    ("API keys not in code",          "OK",
     "OpenRouter key in secrets.dart (template-only today). Never committed.",
     GREEN),
    ("Secrets.dart in .gitignore",    "OK",
     ".gitignore line 47: lib/core/config/secrets.dart. Verified.",
     GREEN),
    ("PB admin creds in CI",          "OK",
     "GitHub Actions secrets only. Not in any yml file. Verified.",
     GREEN),
    ("Barcode cache: anonymous read", "OK",
     "PB listRule='' → world-readable. Only OFF data, no PII.",
     GREEN),
    ("Crash reporting",               "MISSING",
     "No Sentry, no Crashlytics. PII may end up in error logs unfiltered.",
     RED),
    ("TLS pinning",                   "MISSING",
     "Dio talks to PB and OpenRouter over HTTPS, but cert pinning not configured. Acceptable for v1; revisit for v2.",
     YELLOW),
    ("Image storage on device",       "PARTIAL",
     "Camera captures go to app-private dir (good). But they're never deleted; could leak via forensic tools.",
     YELLOW),
    ("PII in error logs",             "RISK",
     "AppErrorHandler captures PlatformDispatcher.onError — should redact email, weight, food names before logging.",
     RED),
    ("Supabase auth",                 "MISSING",
     "Auth is currently local-only. Once you re-add Supabase, must configure RLS on every table from day 1.",
     YELLOW),
    ("Dependency CVE scan",           "MISSING",
     "No automated npm/advisory scan. GitHub Dependabot could enable in 5 minutes.",
     YELLOW),
    ("ProGuard / R8 obfuscation",     "ENABLED",
     "android/app/build.gradle: minifyEnabled true, shrinkResources true. ✓",
     GREEN),
    ("Debug signing for release",     "RISK",
     "build.gradle uses the debug keystore for release builds. Acceptable for early access; MUST replace before public launch.",
     RED),
]
y = Inches(1.55)
for i, (k, status, sub, c) in enumerate(sec):
    add_round(s, Inches(0.5), y + i*Inches(0.43), Inches(12.3), Inches(0.4),
              SURFACE if i % 2 == 0 else SURFACE_MUTE, line=DIVIDER, radius=0.2)
    add_text(s, Inches(0.7), y + i*Inches(0.43) + Inches(0.06), Inches(3.2), Inches(0.3),
             k, size=11, bold=True, color=TEXT_PRI, font='Inter')
    add_text(s, Inches(4.0), y + i*Inches(0.43) + Inches(0.06), Inches(1.4), Inches(0.3),
             status, size=10, bold=True, color=c, font='Inter')
    add_text(s, Inches(5.5), y + i*Inches(0.43) + Inches(0.06), Inches(7.2), Inches(0.3),
             sub, size=9, color=TEXT_SEC, font='Inter', line_spacing=1.2)

page_footer(s, 9, 14)


# ── Slide 10: UX Suggestions ────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
page_bg(s)
page_header(s, "Section 09", "UX Suggestions",
           "What separates \"functional app\" from \"delightful app\"")

ux = [
    ("Onboarding",          "MISSING",
     "3-screen first-run: welcome → why-we're-different → TDEE wizard. Without it users churn in 30s.",
     "2 days", "CRITICAL", RED),
    ("Empty states",        "MISSING",
     "Every screen needs a 'nothing here yet' moment with personality + a CTA. 8 illustrations needed.",
     "1 day", "HIGH", ORANGE),
    ("Haptic feedback",     "PARTIAL",
     "HapticFeedback.mediumImpact on meal log, workout complete, set save. iOS-native feel.",
     "2 hours", "MED", YELLOW),
    ("Sound design",        "MISSING",
     "Swoosh on log, bell on workout complete. just_audio already in pubspec, just unused.",
     "4 hours", "MED", YELLOW),
    ("Streaks + badges",    "MISSING",
     "5-day logging streak, 7-day protein goal hit. Duolingo's #1 retention mechanism. Pure motivation.",
     "3 days", "HIGH", ORANGE),
    ("Pull-to-refresh",     "MISSING",
     "Dashboard re-aggregates. Insights re-queries. Standard mobile pattern.",
     "2 hours", "LOW", GREEN),
    ("Quick actions",       "MISSING",
     "Long-press dashboard tile → quick log. Swipe-to-delete meal. iOS share extension later.",
     "1 day", "MED", YELLOW),
    ("Offline indicator",   "MISSING",
     "Banner: 'Offline — logging locally, will sync when reconnected' + queue count.",
     "3 hours", "MED", YELLOW),
    ("Wow moment",          "MISSING",
     "First time user logs via voice or camera → confetti + 'Saved 2 minutes vs typing!'",
     "3 hours", "HIGH", ORANGE),
    ("Error messages",      "PARTIAL",
     "AppErrorHandler exists but messages should be human, specific, actionable.",
     "1 day", "MED", YELLOW),
    ("a11y",                "MISSING",
     "No Semantics audit, no contrast check, no dynamic font scaling tested.",
     "2 days", "MED", YELLOW),
    ("i18n",                "MISSING",
     "English-only is fine for v1, but structure (.arb) must support it. ~150 strings to externalize.",
     "2-3 days", "LOW", GREEN),
]
y = Inches(1.55)
for i, (k, status, sub, effort, pri, c) in enumerate(ux):
    add_round(s, Inches(0.5), y + i*Inches(0.43), Inches(12.3), Inches(0.4),
              SURFACE if i % 2 == 0 else SURFACE_MUTE, line=DIVIDER, radius=0.2)
    add_text(s, Inches(0.7), y + i*Inches(0.43) + Inches(0.06), Inches(2.0), Inches(0.3),
             k, size=11, bold=True, color=TEXT_PRI, font='Inter')
    add_text(s, Inches(2.7), y + i*Inches(0.43) + Inches(0.06), Inches(1.4), Inches(0.3),
             status, size=10, bold=True, color=c, font='Inter')
    add_text(s, Inches(4.1), y + i*Inches(0.43) + Inches(0.06), Inches(5.7), Inches(0.3),
             sub, size=9, color=TEXT_SEC, font='Inter')
    add_text(s, Inches(9.9), y + i*Inches(0.43) + Inches(0.06), Inches(1.5), Inches(0.3),
             effort, size=10, color=TEXT_PRI, font='Inter', align=PP_ALIGN.CENTER)
    add_text(s, Inches(11.5), y + i*Inches(0.43) + Inches(0.06), Inches(1.3), Inches(0.3),
             pri, size=10, bold=True, color=c, font='Inter', align=PP_ALIGN.CENTER)

page_footer(s, 10, 14)


# ── Slide 11: TOP PROBLEMS — Critical (Now) ─────────────────────────────
s = prs.slides.add_slide(BLANK)
page_bg(s)
page_header(s, "Section 10", "Top Problems · Critical · Act NOW",
           "Blockers that prevent any further work")

# Each card spans full width
problems = [
    ("01", "Build is broken — 31 compile errors",
     "lib/features/barcode/presentation/screens/barcode_scanner_screen.dart (27 errors) + lib/features/camera/presentation/screens/camera_screen.dart (4 errors)",
     "RRect.topLeft/topRight/bottomLeft/bottomRight undefined (use Rect or RRect.tlRadius). 3 callbacks with `Function()` instead of `Function(String)`. MobileScannerErrorBuilder API drift.",
     "App won't run on any device. Nothing else matters until this is fixed.",
     "45 min", "TASK #1: Mechanical patch"),
    ("02", "Zero tests, no test/ directory",
     "No unit tests, no widget tests, no integration tests.",
     "flutter test --coverage in CI passes only because there are no tests to fail. Every future refactor is a coin flip.",
     "Regressions slip through silently. New contributors can't find patterns to copy.",
     "1-2 days", "TASK #5 in roadmap"),
    ("03", "No onboarding, no first-run experience",
     "Initial route is /dashboard. No /login, /signup, or /onboarding route exists.",
     "User opens the app, sees generic 2000 kcal target with no context. Churns in 30 seconds.",
     "Without onboarding, retention is impossible to measure or improve.",
     "2 days", "TASK #6 in roadmap"),
    ("04", "All data wiring is fake / sample",
     "Insights screen renders `_generateSampleWeights()`. Workout screen has no session UI. Dashboard donut takes fake numbers. Camera/barcode review sheets show static UI.",
     "5 of 6 main screens are basically tech demos. Drift tables exist but are unused.",
     "Without real data, the app cannot ship. No amount of polish hides this.",
     "4 days", "TASK #7 in roadmap"),
]
y = Inches(1.55)
for i, (n, title, where, what, impact, effort, fix) in enumerate(problems):
    add_round(s, Inches(0.5), y + i*Inches(1.42), Inches(12.3), Inches(1.32),
              SURFACE, line=DIVIDER, radius=0.05)
    # Number badge
    add_round(s, Inches(0.65), y + i*Inches(1.42) + Inches(0.2), Inches(0.9), Inches(0.9),
              RED, radius=0.1)
    add_text(s, Inches(0.65), y + i*Inches(1.42) + Inches(0.32), Inches(0.9), Inches(0.65),
             n, size=28, bold=True, color=CREAM, font='Inter', align=PP_ALIGN.CENTER)
    # Title
    add_text(s, Inches(1.75), y + i*Inches(1.42) + Inches(0.1), Inches(7.4), Inches(0.3),
             title, size=13, bold=True, color=TEXT_PRI, font='Inter')
    # Where
    add_text(s, Inches(1.75), y + i*Inches(1.42) + Inches(0.38), Inches(7.4), Inches(0.22),
             "WHERE: " + where, size=8, color=TEXT_TER, font='Inter', line_spacing=1.1)
    # What
    add_text(s, Inches(1.75), y + i*Inches(1.42) + Inches(0.6), Inches(7.4), Inches(0.34),
             "WHAT: " + what, size=9, color=TEXT_PRI, font='Inter', line_spacing=1.15)
    # Impact
    add_text(s, Inches(1.75), y + i*Inches(1.42) + Inches(0.95), Inches(7.4), Inches(0.32),
             "IMPACT: " + impact, size=9, bold=True, color=RED, font='Inter', line_spacing=1.15)
    # Effort + fix box
    add_round(s, Inches(9.35), y + i*Inches(1.42) + Inches(0.18), Inches(3.3), Inches(0.95),
              SURFACE_MUTE, line=DIVIDER, radius=0.1)
    add_text(s, Inches(9.5), y + i*Inches(1.42) + Inches(0.28), Inches(3.0), Inches(0.25),
             "EFFORT", size=9, bold=True, color=TEXT_TER, font='Inter')
    add_text(s, Inches(9.5), y + i*Inches(1.42) + Inches(0.5), Inches(3.0), Inches(0.3),
             effort, size=13, bold=True, color=ORANGE, font='Inter')
    add_text(s, Inches(9.5), y + i*Inches(1.42) + Inches(0.82), Inches(3.0), Inches(0.25),
             fix, size=8, color=TEXT_SEC, font='Inter')

page_footer(s, 11, 14)


# ── Slide 12: TOP PROBLEMS — High (this week) ───────────────────────────
s = prs.slides.add_slide(BLANK)
page_bg(s)
page_header(s, "Section 11", "Top Problems · High Priority · This Week",
           "After the critical path, address these")

problems = [
    ("05", "No crash reporting (Sentry/Crashlytics)",
     "No sentry_flutter, no firebase_crashlytics dep. AppErrorHandler captures errors but doesn't report them anywhere.",
     "Production without crash reporting = blind. You can't fix bugs you don't know about. PII risk if errors aren't redacted before logging.",
     "2 hours", "Add sentry_flutter, wrap AppErrorHandler"),
    ("06", "Voice logging UI half-built",
     "record: ^5.1.2 + Whisper path defined in AI gateway. Quick log bar has a mic icon. But no hold-to-talk button, no waveform animation, no transcription feedback.",
     "This is NutriTrack's #1 differentiator vs MyFitnessPal. Without it, app is just another calorie tracker.",
     "3 days", "Hold-to-talk + waveform + Whisper wiring"),
    ("07", "AI camera lacks verification flow",
     "Camera captures, AI returns FoodLogEntry list with confidence. But no UI to show confidence bars or ask user to confirm low-confidence items.",
     "Without verification, users lose trust the moment AI guesses wrong. Need red/yellow/green confidence indicator + 'Take another photo?' escape hatch.",
     "3 days", "Add confidence UI + quick-fix shortcuts"),
    ("08", "Exercise DB: 46 of 488",
     "scripts/exercises_chest.js is the only seed file. Back/shoulders/arms/legs/glutes/core/cardio/full-body missing.",
     "Workout feature is half-built. User can search but most muscle groups show zero results.",
     "4-6 hours", "Boilerplate: 8 more .js files"),
    ("09", "Workout session UI missing",
     "WorkoutScreen has search + list. But no start session → log sets → rest timer → end session flow. No active-session screen, no PR tracking.",
     "Half of NutriTrack's value prop is the workout side. Without session UI, it's a search tool, not a tracker.",
     "4 days", "Active session + set logger + rest timer"),
]
y = Inches(1.55)
for i, (n, title, what, impact, effort, fix) in enumerate(problems):
    add_round(s, Inches(0.5), y + i*Inches(1.1), Inches(12.3), Inches(1.04),
              SURFACE, line=DIVIDER, radius=0.05)
    add_round(s, Inches(0.65), y + i*Inches(1.1) + Inches(0.15), Inches(0.7), Inches(0.7),
              ORANGE, radius=0.1)
    add_text(s, Inches(0.65), y + i*Inches(1.1) + Inches(0.22), Inches(0.7), Inches(0.55),
             n, size=20, bold=True, color=CREAM, font='Inter', align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.55), y + i*Inches(1.1) + Inches(0.08), Inches(7.4), Inches(0.3),
             title, size=12, bold=True, color=TEXT_PRI, font='Inter')
    add_text(s, Inches(1.55), y + i*Inches(1.1) + Inches(0.36), Inches(7.4), Inches(0.34),
             what, size=9, color=TEXT_PRI, font='Inter', line_spacing=1.15)
    add_text(s, Inches(1.55), y + i*Inches(1.1) + Inches(0.7), Inches(7.4), Inches(0.32),
             "IMPACT: " + impact, size=9, bold=True, color=ORANGE, font='Inter', line_spacing=1.15)
    # Effort / fix
    add_round(s, Inches(9.1), y + i*Inches(1.1) + Inches(0.12), Inches(3.55), Inches(0.8),
              SURFACE_MUTE, line=DIVIDER, radius=0.1)
    add_text(s, Inches(9.25), y + i*Inches(1.1) + Inches(0.2), Inches(1.5), Inches(0.22),
             "EFFORT", size=8, bold=True, color=TEXT_TER, font='Inter')
    add_text(s, Inches(9.25), y + i*Inches(1.1) + Inches(0.4), Inches(1.5), Inches(0.3),
             effort, size=12, bold=True, color=ORANGE, font='Inter')
    add_text(s, Inches(10.7), y + i*Inches(1.1) + Inches(0.2), Inches(2.0), Inches(0.22),
             "FIX", size=8, bold=True, color=TEXT_TER, font='Inter')
    add_text(s, Inches(10.7), y + i*Inches(1.1) + Inches(0.4), Inches(2.0), Inches(0.4),
             fix, size=8, color=TEXT_PRI, font='Inter', line_spacing=1.15)

page_footer(s, 12, 14)


# ── Slide 13: TOP PROBLEMS — Medium (this month) ────────────────────────
s = prs.slides.add_slide(BLANK)
page_bg(s)
page_header(s, "Section 12", "Top Problems · Medium Priority · This Month",
           "Polish, performance, security hardening")

problems = [
    ("10", "9 packages on major N-1",
     "flutter_riverpod 2.6 (latest 3.3), freezed 2.5 (3.2), mobile_scanner 5.2 (7.2), health REMOVED, go_router 14.8 (17.3), intl 0.19 (0.20), etc.",
     "Future Flutter versions break these. Each is a real migration: API renames, codegen output changes, permission flow drift.",
     "1 week", "5 grouped PRs per KNOWN_WARNINGS §6"),
    ("11", "7 dead dependencies in pubspec",
     "just_audio, web_socket_channel, confetti, flutter_staggered_animations, animations, lottie, gap, shimmer, supabase_flutter.",
     "0 imports in lib/ for most. ~10 MB wasted APK size. Slower builds.",
     "30 min", "flutter pub remove for each"),
    ("12", "README lies about Isar",
     "Stack table says Local DB | Isar 3.1. You migrated to Drift 5 commits ago.",
     "Confuses every new contributor. Costs an hour of confusion per onboarding.",
     "10 min", "3-line edit + bump version"),
    ("13", "Release build uses debug keystore",
     "android/app/build.gradle: signingConfigs.release = signingConfigs.debug",
     "Acceptable for internal testing. CRITICAL to replace before public launch — Play Store / App Store will reject.",
     "1 hour", "Generate release keystore, store in key.properties"),
    ("14", "Adaptive TDEE engine not built",
     "TDEE wizard calculates starting value. But no weekly recalibration from real weight + intake data.",
     "MacroFactor's killer feature. Without it, NutriTrack's target is static even after 6 months of logging.",
     "1 day", "7700 kcal/kg formula + weekly trigger"),
    ("15", "Image storage on device",
     "Camera captures saved to app-private dir. Never deleted.",
     "Old food photos pile up. Forensic tools can recover. PII risk.",
     "2 hours", "Retention policy + cleanup cron"),
]
y = Inches(1.55)
for i, (n, title, what, impact, effort, fix) in enumerate(problems):
    add_round(s, Inches(0.5), y + i*Inches(0.93), Inches(12.3), Inches(0.86),
              SURFACE, line=DIVIDER, radius=0.05)
    add_round(s, Inches(0.65), y + i*Inches(0.93) + Inches(0.15), Inches(0.55), Inches(0.55),
              YELLOW, radius=0.1)
    add_text(s, Inches(0.65), y + i*Inches(0.93) + Inches(0.2), Inches(0.55), Inches(0.45),
             n, size=14, bold=True, color=TEXT_PRI, font='Inter', align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.4), y + i*Inches(0.93) + Inches(0.06), Inches(7.5), Inches(0.28),
             title, size=12, bold=True, color=TEXT_PRI, font='Inter')
    add_text(s, Inches(1.4), y + i*Inches(0.93) + Inches(0.3), Inches(7.5), Inches(0.34),
             what, size=9, color=TEXT_SEC, font='Inter', line_spacing=1.15)
    add_text(s, Inches(1.4), y + i*Inches(0.93) + Inches(0.6), Inches(7.5), Inches(0.26),
             "IMPACT: " + impact, size=9, bold=True, color=YELLOW, font='Inter', line_spacing=1.15)
    add_round(s, Inches(9.1), y + i*Inches(0.93) + Inches(0.13), Inches(3.55), Inches(0.6),
              SURFACE_MUTE, line=DIVIDER, radius=0.1)
    add_text(s, Inches(9.25), y + i*Inches(0.93) + Inches(0.18), Inches(1.5), Inches(0.22),
             "EFFORT", size=8, bold=True, color=TEXT_TER, font='Inter')
    add_text(s, Inches(9.25), y + i*Inches(0.93) + Inches(0.38), Inches(1.5), Inches(0.25),
             effort, size=11, bold=True, color=ORANGE, font='Inter')
    add_text(s, Inches(10.7), y + i*Inches(0.93) + Inches(0.18), Inches(2.0), Inches(0.22),
             "FIX", size=8, bold=True, color=TEXT_TER, font='Inter')
    add_text(s, Inches(10.7), y + i*Inches(0.93) + Inches(0.38), Inches(2.0), Inches(0.4),
             fix, size=8, color=TEXT_PRI, font='Inter', line_spacing=1.15)

page_footer(s, 13, 14)


# ── Slide 14: Next Steps ────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
page_bg(s)
# Big orange background
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, ORANGE)
# Cream panel
add_round(s, Inches(0.7), Inches(0.7), Inches(11.9), Inches(6.1), CREAM, radius=0.02)

add_text(s, Inches(1.0), Inches(1.1), Inches(10), Inches(0.4),
         "WHAT TO DO NEXT", size=12, bold=True, color=ORANGE, font='Inter')
add_text(s, Inches(1.0), Inches(1.45), Inches(11), Inches(1.0),
         "Three steps to get\nNutriTrack shipping.", size=36, bold=True,
         color=TEXT_PRI, font='Inter', line_spacing=1.1)

steps = [
    ("TODAY",
     "Fix the build.",
     "Run the 5-line edit list for RRect + Function() signatures. 45 minutes. Open the icon. Run flutter run."),
    ("THIS WEEK",
     "Wire the data + add Sentry.",
     "Replace 5 sample-data screens with real Drift. Add Sentry (2 hours). Wire 7 priority tests. ~3 days."),
    ("THIS MONTH",
     "Ship differentiators.",
     "Voice UI + camera verification + workout session + complete exercise DB. ~3 weeks. Then TestFlight + public beta."),
]
y = Inches(3.0); h = Inches(1.15); gap = Inches(0.15)
for i, (when, what, how) in enumerate(steps):
    add_round(s, Inches(1.0), y + i*(h+gap), Inches(11.0), h, SURFACE, line=DIVIDER, radius=0.06)
    add_rect(s, Inches(1.0), y + i*(h+gap), Inches(0.12), h, ORANGE)
    add_text(s, Inches(1.3), y + i*(h+gap) + Inches(0.12), Inches(2.5), Inches(0.4),
             when, size=12, bold=True, color=ORANGE, font='Inter')
    add_text(s, Inches(1.3), y + i*(h+gap) + Inches(0.42), Inches(3.5), Inches(0.5),
             what, size=18, bold=True, color=TEXT_PRI, font='Inter')
    add_text(s, Inches(5.0), y + i*(h+gap) + Inches(0.32), Inches(6.8), Inches(0.7),
             how, size=12, color=TEXT_SEC, font='Inter', line_spacing=1.4)

add_text(s, Inches(1.0), Inches(6.55), Inches(11), Inches(0.3),
         "Repo: github.com/albertlaudia/nutritrack  ·  PB: pocketbase.scaleupcrm.com  ·  9,373 Dart LOC",
         size=10, color=TEXT_TER, font='Inter')

page_footer(s, 14, 14)

# ── Save ────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Size:  {os.path.getsize(OUT):,} bytes")
print(f"Slides: {len(prs.slides)}")